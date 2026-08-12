"""발표자료(=보고서) PPT 생성 -> submission/02_발표자료/분석보고서_문지영.pptx

장수를 정한 근거 (submission/02_발표자료/발표가이드.png)
--------------------------------------------------
발표는 **10분이며 초과하면 강제 종료**된다. "세부 내용보다는 핵심 분석 결과와
결론 중심으로" 라고 명시돼 있다. 가이드의 권장 배분이 곧 섹션 설계다.

  주제·목적·핵심질문   1분      대시보드 시연   2분
  데이터·이슈·해결     2분      한계와 배운 점  1분
  분석 과정·결과       4분

시연 2분을 빼면 슬라이드는 8분이다. 간지 4장을 빼면 내용은 21장이라
장당 23초 안팎으로 잡았다. 처음 계획한 25장(장당 19초)은 지킬 수 없는 분량이었다.

디자인 — 예시 덱(SKN21 2팀)에서 따온 규칙
--------------------------------------
· 16:9, 흰 배경. 색 블록으로 화면을 채우지 않고 **여백**으로 나눈다
· 좌상단 **파란 알약 배지**에 섹션명, 그 아래 큰 제목 (제목 줄을 따로 그리지 않는다)
· 본문은 불릿 2~3개, 강조는 굵게. 문장을 길게 늘어놓지 않는다
· 근거 그림은 아래쪽에 크게. 출처·단위는 오른쪽 아래 작은 회색
· 쪽번호는 오른쪽 아래
· 표지·간지는 큰 활자 + 얇은 가로선 + 하단 2단 푸터

그림
----
`src/report/deck_figures.py`가 만든 슬라이드 전용 그림을 쓴다.
분석용 그림(reports/figures/*.png)은 패널이 3~4개라 슬라이드에서 읽히지 않는다.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.utils.logger import get_logger
from src.utils.settings import ROOT_DIR

logger = get_logger(__name__)

DECK_FIG = ROOT_DIR / "reports" / "figures" / "deck"
OUT_PATH = ROOT_DIR / "submission" / "02_발표자료" / "분석보고서_문지영.pptx"

# 예시 덱과 같은 16:9
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

FONT = "맑은 고딕"
INK = RGBColor(0x11, 0x11, 0x11)
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x8A, 0x8A, 0x8A)
LINE = RGBColor(0xD8, 0xD8, 0xD8)
BLUE = RGBColor(0x18, 0x77, 0xF2)
RED = RGBColor(0xD0, 0x3B, 0x3B)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

M_LEFT = 0.9                 # 좌측 기준선 — 배지·제목·본문이 모두 여기에 맞는다
TITLE_Y = 1.15
BODY_Y = 2.25

SUBTITLE = "생활인구·거주형태 기반 서울 공영주차 수급 진단"
AUTHOR = "문지영"
DATE = "2026.08."


# ── 저수준 헬퍼 ────────────────────────────────────────────────
def _tf(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    return box.text_frame


def _split_lines(runs):
    """조각 안의 '\\n'을 **문단 분리**로 바꾼다.

    python-pptx는 run 텍스트의 개행을 `<a:br/>`로 만들지 않고 `<a:t>` 안에
    그대로 넣는다. PowerPoint는 이를 줄바꿈으로 그려주지만 규격상 보장이 아니라,
    뷰어에 따라 공백 하나로 붙어버린다. 문단으로 쪼개면 어디서나 줄이 나뉜다.
    """
    out, cur = [], []
    for item in runs:
        pieces = item if isinstance(item, list) else [item]
        for piece in pieces:
            text, opt = (piece, {}) if isinstance(piece, str) else piece
            for k, part in enumerate(str(text).split("\n")):
                if k > 0:
                    out.append(cur)
                    cur = []
                cur.append((part, opt))
        out.append(cur)
        cur = []
    return out


def _write(frame, runs, *, size=16, color=BODY, bold=False, align=PP_ALIGN.LEFT,
           space=6, line=1.35):
    """runs = [문자열 | (문자열, {옵션})]. 옵션은 size/bold/color.

    줄간격을 배수가 아니라 **절대값(포인트)** 으로 넣는다.
    python-pptx에 실수를 주면 `<a:lnSpc><a:spcPct val="135000"/>`(135%)로 나가는데,
    PowerPoint는 이를 제대로 읽지만 일부 뷰어(VS Code의 Muty PPT Viewer 등)는
    배수를 무시해 **여러 줄이 한 자리에 겹쳐 찍힌다.**
    포인트로 주면 `<a:spcPts>`가 되어 어느 렌더러에서도 같은 높이가 나온다.
    """
    for i, item in enumerate(_split_lines(runs)):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        # 한 문단 안에서 굵기를 섞으려면 (조각, 옵션) 리스트를 준다
        pieces = item if isinstance(item, list) else [item]
        sizes = [(piece[1].get("size", size) if isinstance(piece, tuple) else size)
                 for piece in pieces]
        p.line_spacing = Pt(max(sizes) * line)
        for piece in pieces:
            text, opt = (piece, {}) if isinstance(piece, str) else piece
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.color.rgb = opt.get("color", color)


def _pill(slide, x, y, text, *, fill=BLUE, color=WHITE, size=12, w=None):
    """알약 배지. 글자 길이에 맞춰 폭을 잡는다."""
    w = w or (0.34 + len(text) * (size / 72) * 1.15)
    h = 0.36
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = 0.5           # 완전한 알약 모양
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write(tf, [text], size=size, bold=True, color=color, align=PP_ALIGN.CENTER, space=0, line=1.0)
    return s


def _hairline(slide, y, *, x=M_LEFT, w=None):
    w = w or (13.333 - 2 * M_LEFT)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = LINE
    s.line.fill.background()
    s.shadow.inherit = False


def _page_no(slide, n):
    _write(_tf(slide, 12.2, 6.82, 0.8, 0.4), [str(n)],
           size=12, color=MUTED, align=PP_ALIGN.RIGHT, space=0)


def _note(slide, text, *, y=6.62):
    """출처·단위 같은 잔글씨. 오른쪽 아래 — 쪽번호(12.2")를 침범하지 않게 폭을 잡는다."""
    _write(_tf(slide, 3.6, y, 8.4, 0.35), [text],
           size=10.5, color=MUTED, align=PP_ALIGN.RIGHT, space=0)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _content(prs, n, section, title, *, sub=None):
    """내용 슬라이드의 머리 — 배지 + 제목 (+ 한 줄 설명)."""
    s = _blank(prs)
    _pill(s, M_LEFT, 0.55, section)
    _write(_tf(s, M_LEFT - 0.06, TITLE_Y, 11.5, 0.8), [title],
           size=30, bold=True, color=INK, space=0)
    if sub:
        _write(_tf(s, M_LEFT - 0.06, TITLE_Y + 0.72, 11.5, 0.5), [sub],
               size=14, color=MUTED, space=0)
    _page_no(s, n)
    return s


def _bullets(slide, items, *, y=BODY_Y, size=15.5, w=11.5):
    """예시 덱처럼 '•' + 굵은 강조. items = [문자열 | 조각 리스트]

    상자 높이를 줄 수에 맞춰 잡는다. 예전에는 1.6"로 고정해 두어
    뒤에 오는 주석·쪽번호와 겹치고 슬라이드 밖으로 나가기도 했다.
    """
    runs = []
    for it in items:
        pieces = it if isinstance(it, list) else [it]
        runs.append([("•  ", {"color": BLUE, "bold": True})] + list(pieces))
    h = 0.4 * len(items) + 0.1
    _write(_tf(slide, M_LEFT, y, w, h), runs, size=size, color=BODY, space=7)
    return y + h


def _picture(slide, name, *, top, height, width=11.6):
    path = DECK_FIG / f"{name}.png"
    if not path.exists():
        logger.warning(f"그림 없음: {name}")
        return
    pic = slide.shapes.add_picture(str(path), 0, Inches(top))
    scale = min(Inches(width) / pic.width, Inches(height) / pic.height)
    pic.width, pic.height = int(pic.width * scale), int(pic.height * scale)
    pic.left = int((SLIDE_W - pic.width) / 2)


def _table(slide, rows, *, y, col_w, size=13, x=M_LEFT, w=11.5, row_h=0.42):
    """표를 그리고 **아래쪽 y 좌표**를 돌려준다.

    PowerPoint는 셀 글자가 넘치면 행 높이를 자동으로 늘린다.
    그래서 `row_h × 행수`로 다음 요소를 배치하면 표가 그 위를 덮는다.
    행 높이를 명시해 최소값을 고정하고, 호출부는 반환값 기준으로 다음 요소를 놓는다.
    (셀 글자가 한 줄에 들어가도록 열 폭·글자 크기를 잡는 것이 전제다)
    """
    n_r, n_c = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(row_h * n_r))
    tbl = shape.table
    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(w * cw / total)
    for r in tbl.rows:
        r.height = Inches(row_h)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.12)
            head = ri == 0
            _write(cell.text_frame, [str(val)], size=size, bold=head,
                   color=INK if head else BODY, space=0, line=1.15)
            cell.fill.solid()
            # 머리글만 옅은 회색, 본문은 흰색 — 예시 덱처럼 표에 색을 얹지 않는다
            cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7) if head else WHITE
    return y + row_h * n_r


def _stat(slide, x, y, value, label, *, color=BLUE, w=2.6, size=40):
    """큰 숫자 한 덩어리. 값이 길면 size를 줄여 줄바꿈으로 라벨을 밀지 않게 한다."""
    _write(_tf(slide, x, y, w, 0.85), [value], size=size, bold=True, color=color,
           align=PP_ALIGN.CENTER, space=0)
    _write(_tf(slide, x, y + 0.9, w, 0.5), [label], size=13, color=MUTED,
           align=PP_ALIGN.CENTER, space=0)


# ── 표지 · 간지 ────────────────────────────────────────────────
def _footer_block(slide):
    _hairline(slide, 5.95)
    _write(_tf(slide, M_LEFT, 6.15, 7.0, 0.9),
           [[(SUBTITLE, {"bold": True, "size": 14, "color": INK})],
            [(AUTHOR, {"bold": True, "size": 14, "color": INK}),
             ("    SKN 데이터분석 · 1주 개인 프로젝트", {"size": 13, "color": MUTED})]],
           space=3)
    _write(_tf(slide, 8.0, 6.15, 4.4, 0.9),
           [DATE, "서울 424개 행정동"],
           size=13.5, color=INK, align=PP_ALIGN.RIGHT, space=3)


def s_cover(prs, n):
    s = _blank(prs)
    _write(_tf(s, M_LEFT, 1.35, 11.5, 0.6),
           [("🅿  공영주차장은 사람이 오는 곳에 있는가?", {"size": 20, "bold": True, "color": INK})],
           space=0)
    # 큰 제목은 줄마다 별도 상자에 둔다.
    # 한 상자에 두 문단으로 넣으면 줄간격을 무시하는 뷰어에서 두 줄이 겹쳐 찍힌다.
    _write(_tf(s, M_LEFT - 0.12, 2.25, 11.8, 1.15),
           [("서울 공영주차", {"size": 62, "bold": True, "color": INK})], space=0)
    _write(_tf(s, M_LEFT - 0.12, 3.35, 11.8, 1.15),
           [("수급 진단", {"size": 62, "bold": True, "color": BLUE})], space=0)
    _footer_block(s)
    _page_no(s, n)


def s_agenda(prs, n):
    s = _content(prs, n, "목차", "Contents")
    rows = [
        ("01", "분석 배경 및 목표", "왜 이 분석이 필요한가"),
        ("02", "활용 데이터", "무엇을 어떻게 붙였는가"),
        ("03", "분석 결과", "가설 6개의 판정"),
        ("04", "결론", "두 사용자에게 무엇을 주는가"),
    ]
    y = 2.5
    for no, title, desc in rows:
        _write(_tf(s, M_LEFT, y, 1.2, 0.6), [no], size=26, bold=True, color=BLUE, space=0)
        _write(_tf(s, M_LEFT + 1.25, y + 0.04, 9.5, 0.6),
               [[(title, {"size": 22, "bold": True, "color": INK}),
                 (f"     {desc}", {"size": 14, "color": MUTED})]], space=0)
        y += 0.95


def s_divider(prs, n, no, title):
    s = _blank(prs)
    _write(_tf(s, M_LEFT, 2.05, 4.0, 0.9), [no], size=40, bold=True, color=BLUE, space=0)
    _write(_tf(s, M_LEFT - 0.12, 2.85, 11.8, 1.4), [title],
           size=54, bold=True, color=INK, space=0)
    _footer_block(s)
    _page_no(s, n)


# ── I. 분석 배경 및 목표 ───────────────────────────────────────
def s_scope(prs, n):
    """발표가이드가 첫 1분에 요구하는 '주제·목적·핵심 질문'을 한 장에 담는다.

    목적은 '왜 하는가'(두 사용자별), 목표는 '무엇을 만들어 내는가'를
    **숫자로 확인 가능한 형태**로 적는다. 뒤 슬라이드가 이 숫자를 하나씩 채운다.
    """
    s = _content(prs, n, "분석 배경 및 목표", "주제 · 목적 · 목표")

    def row(y, label, runs, *, size=15):
        _pill(s, M_LEFT, y + 0.02, label, size=12.5, w=0.95)
        _write(_tf(s, M_LEFT + 1.2, y - 0.04, 10.9, 1.8), runs, size=size, space=6)

    row(2.15, "주제",
        [[("서울 424개 행정동의 공영주차 수급을 ", {"size": 18}),
          ("생활인구·거주형태", {"size": 18, "bold": True, "color": BLUE}),
          ("로 진단하고", {"size": 18})],
         [("확충 우선순위를 도출한다", {"size": 18})]])

    row(3.35, "목적",
        [[("나들이객 ", {"color": MUTED}),
          ("목적지를 정하기 이전 단계의 의사결정을 돕는다", {"color": INK})],
         [("정책 담당 ", {"color": MUTED}),
          ("확충 우선순위를 판단할 정량 근거를 만든다", {"color": INK})]])

    row(4.55, "목표",
        [[("행정동 424 × 요일 7 × 시간대 5 = ", {}),
          ("14,840행 패널", {"bold": True, "color": INK}), (" 구축", {})],
         [("검증 가설 ", {}), ("6개", {"bold": True, "color": INK}),
          (" — 유동수요·정주수요 양쪽으로 공급을 설명해 본다", {})],
         [("추천 동네 ", {}), ("23개", {"bold": True, "color": INK}),
          (" · 확충 후보 ", {}), ("52개", {"bold": True, "color": INK}),
          (" · 확충 불필요 ", {}), ("13개", {"bold": True, "color": INK}), (" 도출", {})],
         [("두 사용자가 각자 쓰는 ", {}),
          ("대시보드 5페이지", {"bold": True, "color": INK})]])

    _note(s, "핵심 질문 — \"공영주차장은 사람이 오는 곳에 있는가?\"", y=6.6)


def s_problem(prs, n):
    s = _content(prs, n, "분석 배경 및 목표", "문제 정의",
                 sub="기존 서비스는 목적지가 이미 정해진 사람만 돕고, 행정에는 확충 우선순위를 판단할 근거가 없다")
    _pill(s, M_LEFT, 2.35, "기존 서비스", fill=RGBColor(0x6B, 0x6B, 0x6B), size=13)
    _write(_tf(s, M_LEFT, 2.85, 5.3, 1.2),
           [[("모두의주차장 · 서울주차정보", {"size": 15, "color": INK})],
            [("\"강남역 근처 주차장 어디?\"", {"size": 17, "bold": True, "color": BODY})]],
           space=5)
    # 오른쪽 칸에 두 사용자의 질문을 나란히 둔다.
    # 예전에는 나들이객의 질문만 있어 부 사용자가 아래 표에서 갑자기 나왔다.
    _pill(s, 7.0, 2.35, "이 분석이 답하는 것 — 두 사용자", fill=BLUE, size=13)
    _write(_tf(s, 7.0, 2.85, 5.6, 1.3),
           [[("나들이객 ", {"size": 13, "color": MUTED}),
             ("\"어디로 갈까? 언제 갈까?\"", {"size": 16, "bold": True, "color": BLUE})],
            [("정책 담당 ", {"size": 13, "color": MUTED}),
             ("\"어디에 지어야 하나?\"", {"size": 16, "bold": True, "color": BLUE})]],
           space=5)
    _hairline(s, 4.3)
    _write(_tf(s, M_LEFT, 4.5, 11.5, 0.5),
           [[("두 사용자가 ", {"size": 18, "bold": True, "color": INK}),
             ("같은 분석의 양쪽 끝", {"size": 18, "bold": True, "color": BLUE}),
             ("을 쓴다 — 공급이 여유로운 쪽은 \"여기로 가세요\", 부족한 쪽은 \"여기에 지으세요\"",
              {"size": 14, "color": BODY})]],
           space=0)
    _table(s, [
        ["", "묻는 것", "받는 것"],
        ["주 사용자  나들이객", "어느 동네에 언제 가면 주차가 덜 막힐까", "추천 동네 · 동네별 최적 시점"],
        ["부 사용자  자치구 주차 정책 담당", "공영주차장을 어디에 지어야 하나", "확충 후보 · 확충 불필요 지역"],
    ], y=5.1, col_w=[3.0, 4.6, 3.9], size=13)


def s_hypotheses(prs, n):
    s = _content(prs, n, "분석 배경 및 목표", "검증한 가설 6개",
                 sub="가장 중요한 발견은 가설이 기각된 것이다")
    _table(s, [
        ["#", "가설", "판정"],
        ["1", "유동인구가 많은 동일수록 1인당 공영주차면이 낮다", "약한 지지   r = -0.14"],
        ["2", "동네 유형에 따라 시간대 패턴이 다르다", "지지"],
        ["3", "20~30대 비중·음식점 밀도와 주차면에 음의 상관", "기각   r = -0.02"],
        ["4", "매력도와 주차 여유를 동시에 갖춘 예외가 있다", "지지   23개 동"],
        ["5", "시간대 조정만으로 혼잡도가 낮아진다", "지지   +34.5%"],
        ["6", "아파트 거주 비율이 취약지역 오탐을 만든다", "지지   4개 동"],
    ], y=2.45, col_w=[0.5, 7.6, 3.4], size=14, row_h=0.48)
    _bullets(s, [
        [("가설 1·3이 기각된 것이 이 분석의 출발점 — ", {}),
         ("공급이 수요를 따라가지 않는다", {"bold": True, "color": INK})],
        [("가설 6은 ", {}), ("'생활인구만으로 대상지를 고르면 안 된다'", {"bold": True, "color": INK}),
         ("는 실무적 결론으로 이어진다", {})],
    ], y=5.85)


def s_pipeline(prs, n):
    s = _content(prs, n, "분석 배경 및 목표", "분석 과정과 도구",
                 sub="수집부터 산출까지 전 과정을 스크립트로 재현할 수 있다")
    steps = [("수집", "OpenAPI 4종\n수동 CSV 3종"), ("전처리", "코드 크로스워크\n지오코딩 · 패널화"),
             ("분석", "회귀 잔차 · PCA\n군집 · 프리드먼"), ("산출", "보고서 · 대시보드\n확충 우선순위")]
    x = M_LEFT
    for i, (title, desc) in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4),
                                 Inches(2.5), Inches(1.35))
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE if i == 3 else RGBColor(0xF2, 0xF4, 0xF7)
        box.line.color.rgb = BLUE if i == 3 else LINE
        box.shadow.inherit = False
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _write(tf, [[(title, {"size": 17, "bold": True,
                              "color": WHITE if i == 3 else INK})],
                    [(desc, {"size": 11.5, "color": WHITE if i == 3 else MUTED})]],
               align=PP_ALIGN.CENTER, space=2)
        if i < 3:
            _write(_tf(s, x + 2.55, 2.85, 0.4, 0.5), ["›"],
                   size=22, color=MUTED, align=PP_ALIGN.CENTER, space=0)
        x += 2.95
    _write(_tf(s, M_LEFT, 4.15, 11.5, 0.5),
           [("가장 어려웠던 문제", {"size": 18, "bold": True, "color": INK})], space=0)
    _table(s, [
        ["문제", "해결", "결과"],
        ["데이터마다 행정동 코드 체계가 달랐다", "'전국 행정동 코드정보'를 다리로 사용", "조인율 7.5% → 100%"],
        ["주차장에 행정동 정보가 없고 좌표도 결측", "SGIS 지오코딩으로 주소를 코드로 변환", "배정률 100%"],
        ["서울시 API에 1면짜리 거주자우선이 섞임", "표준데이터로 교체 후 0면 동만 보충", "78,267면 / 358개 동"],
    ], y=4.7, col_w=[4.3, 4.4, 2.8], size=12.5)
    _note(s, "Python 3.14 · pandas · scipy · scikit-learn · geopandas · matplotlib · Streamlit · python-pptx")


# ── II. 활용 데이터 ────────────────────────────────────────────
def s_data(prs, n):
    s = _content(prs, n, "활용 데이터", "활용 데이터 7종",
                 sub="공공데이터 6종 + 수동 다운로드 CSV")
    _table(s, [
        ["데이터", "출처", "범위"],
        ["우리마을 생활인구(행정동)", "서울 열린데이터광장", "56일 · 157MB"],
        ["전국주차장 표준데이터", "공공데이터포털", "856개소 73,796면"],
        ["서울시 공영주차장", "서울 열린데이터광장", "보충 43개 동 4,471면"],
        ["상권분석 4종", "서울 열린데이터광장", "점포·집객시설·상주·직장"],
        ["인구주택총조사 거처종류별 가구", "KOSIS", "서울 427개 동"],
        ["행정동 경계", "SGIS", "426개 폴리곤"],
        ["전국 행정동 코드정보", "서울 열린데이터광장", "43개월분"],
    ], y=2.4, col_w=[5.0, 3.4, 3.6], size=13, row_h=0.40)
    _bullets(s, [
        [("주차장을 2종 쓴 이유 — 서울시 API에는 ", {}),
         ("1면짜리 거주자우선 구획", {"bold": True, "color": INK}),
         ("이 섞여 있다 (서교동 6면 vs 152면)", {})],
        [("표준데이터를 기본으로 하고 0면인 동에만 보충해 중복 없이 ", {}),
         ("78,267면", {"bold": True, "color": INK}), ("을 확보했다", {})],
    ], y=5.75)


def s_join(prs, n):
    s = _content(prs, n, "활용 데이터", "데이터를 붙이는 것이 절반이었다",
                 sub="서로 다른 출처를 하나의 행정동 축에 모으는 데 가장 많은 시간이 들었다")
    _stat(s, 1.0, 2.6, "7.5% → 100%", "행정동 코드 조인율", w=3.4, size=34)
    _stat(s, 5.0, 2.6, "100%", "주차장 행정동 배정률", w=3.4)
    _stat(s, 9.0, 2.6, "424 / 424", "총조사 매칭", w=3.4, size=36)
    _hairline(s, 4.35)
    _bullets(s, [
        [("생활인구는 행안부 코드, 경계는 SGIS 코드 — 그냥 붙이면 ", {}),
         ("424개 중 32개만 일치", {"bold": True, "color": RED}),
         ("했다", {})],
        [("'전국 행정동 코드정보' 43개월분을 다리로 놓아 ", {}),
         ("인구 손실 0명", {"bold": True, "color": INK}), ("으로 100% 연결했다", {})],
        [("주차장은 좌표가 5.7~33.5% 결측이라 공간조인 대신 ", {}),
         ("SGIS 지오코딩", {"bold": True, "color": INK}),
         ("으로 주소를 행정동코드로 직접 변환했다", {})],
    ], y=4.6)
    _note(s, "예외 처리 — 가운뎃점 표기 · 구버전 코드 7개 · 동 분할 1:2 · 동명 중복(신사동) · 통합동(용신동)")


def s_timeslot_design(prs, n):
    s = _content(prs, n, "활용 데이터", "시간대를 다시 나눈 이유",
                 sub="처음 4구간으로 나눴더니 실행할 수 없는 추천이 나왔다")
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M_LEFT), Inches(2.35),
                             Inches(11.5), Inches(1.0))
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xF3)
    box.line.color.rgb = RGBColor(0xF0, 0xD5, 0xD5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    _write(tf, [[("\"일요일 아침이 가장 여유롭다\"", {"size": 18, "bold": True, "color": RED}),
                 ("     ← 나들이 수요가 없는 06~08시가 '아침(06-11)' 평균에 섞인 결과",
                  {"size": 13, "color": BODY})]], space=0)
    _table(s, [
        ["", "구간", "구간 내부 편차"],
        ["처음", "아침 06-11 / 점심 / 오후 / 저녁밤 18-24", "1.110"],
        ["바꾼 뒤", "오전 10-12 / 점심 12-15 / 오후 15-18 / 저녁 18-21 / 밤 21-24", "1.043"],
    ], y=3.65, col_w=[1.3, 7.6, 2.6], size=13)
    _write(_tf(s, M_LEFT, 5.05, 11.5, 0.42),
           [("공휴일 2일도 요일 평균에서 제외했다", {"size": 17, "bold": True, "color": INK})],
           space=0)
    _table(s, [
        ["제외한 날", "정체", "효과 — 요일별 변동계수"],
        ["2026-06-03 (수)", "제9회 전국동시지방선거 (임시공휴일)", "수  2.01% → 0.19%"],
        ["2026-07-17 (금)", "제헌절 (2026년 부활)", "금  2.97% → 0.38%"],
    ], y=5.6, col_w=[2.6, 5.6, 3.3], size=12.5)


# ── III. 분석 결과 ─────────────────────────────────────────────
def s_eda(prs, n):
    s = _content(prs, n, "분석 결과", "먼저 데이터를 훑었다 (EDA)",
                 sub="가설을 검정하기 전에 분포·결측·이상치·패턴부터 확인했다")
    _picture(s, "08_eda_corr", top=2.35, height=3.3)
    # 불릿은 한 줄에 들어가야 한다 — 넘치면 상자가 커져 아래 주석과 겹친다
    _bullets(s, [
        [("매력도 변수끼리는 0.75~0.81로 묶이는데 ", {}),
         ("천명당주차면만 -0.17~-0.02", {"bold": True, "color": RED}),
         (" — 가설 1·3 기각의 근거", {})],
        [("이상치는 지우지 않았다 — 서교동 117,158명은 ", {}),
         ("오류가 아니라 분석 대상", {"bold": True, "color": INK}), ("이기 때문", {})],
    ], y=5.7)
    _note(s, "EDA가 바꾼 결정 — 공휴일 제외 · 시간대 5구간 · 0면 동 보충 · 민감도 병기", y=6.75)


def s_core(prs, n):
    s = _content(prs, n, "분석 결과", "공급은 수요를 따라가지 않는다",
                 sub="유동수요와 정주수요 어느 쪽으로도 공영주차 공급이 설명되지 않는다")
    _picture(s, "01_core_result", top=2.45, height=3.25)
    _bullets(s, [
        [("논현2동은 4.4만 명에 13면, 도봉1동은 1.75만 명에 1,413면 — ", {}),
         ("100배 넘는 격차", {"bold": True, "color": RED})],
        [("설명력이 낮다는 것 자체가 발견 — ", {}),
         ("현재 배치가 수요 기반이 아니라는 근거", {"bold": True, "color": INK}), ("다", {})],
    ], y=5.85)
    _note(s, "토요일 오후 · 로그-로그 단순회귀 · 공영주차 0면 66개 제외 "
             "(정주수요는 비아파트 100가구 미만 동도 제외 — 실수요 비교라는 의미를 지키기 위해)")


def s_hyp13(prs, n):
    s = _content(prs, n, "분석 결과", "가설 1·3 — 매력도와 공급의 관계",
                 sub="상관은 방향만 맞고 크기가 없다")
    _picture(s, "02_hypothesis_13", top=2.5, height=2.95)
    _bullets(s, [
        [("가설 3(20~30대·음식점 밀도)은 ", {}), ("기각", {"bold": True, "color": RED}),
         (" — 음식점이 많다고 주차가 부족하지는 않다", {})],
        [("0면 66개 동을 포함하면 음식점 상관은 +0.027로 ", {}),
         ("부호가 뒤집혀", {"bold": True, "color": INK}), (" 두 값을 함께 제시했다", {})],
    ], y=5.7)


def s_cluster(prs, n):
    s = _content(prs, n, "분석 결과", "가설 2 — 동네 유형에 따라 패턴이 뒤집힌다",
                 sub="실루엣 계수 기준 최적 군집 수는 2 (예상한 3유형은 데이터가 지지하지 않음)")
    _picture(s, "03_cluster_pattern", top=2.5, height=3.1)
    _bullets(s, [
        [("상권형 112개 동은 주말·밤에 여유롭고, 주거형 239개 동은 ", {}),
         ("정확히 반대", {"bold": True, "color": INK}), ("다", {})],
        [("같은 '주차 여유'라도 언제 여유로운지가 유형별로 달라 안내 문구를 나눠야 한다", {})],
    ], y=5.8)


def s_timeslot(prs, n):
    s = _content(prs, n, "분석 결과", "가설 5 — 목적지를 바꾸지 않아도 된다",
                 sub="행정동을 블록으로 한 프리드먼 검정  χ² = 1,188.0,  p = 1.3e-227")
    _picture(s, "04_timeslot_effect", top=2.5, height=3.05)
    _bullets(s, [
        [("중구 소공동은 수·점심 대비 일·밤이 ", {}), ("989% 여유", {"bold": True, "color": INK}),
         (" — 상위권이 전부 업무지구라 가설 2와 일관", {})],
        [("공영주차면 수는 변하지 않는다. '주차장이 느는 것'이 아니라 ", {}),
         ("'경쟁할 사람이 주는 것'", {"bold": True, "color": INK}), ("이다", {})],
    ], y=5.75)


def s_recommend(prs, n):
    s = _content(prs, n, "분석 결과", "가설 4 — 놀거리와 주차를 동시에 갖춘 곳",
                 sub="매력도는 PCA로 뽑은 축이며 임의 가중치를 쓰지 않았다")
    _picture(s, "05_recommend", top=2.5, height=3.05)
    _bullets(s, [
        [("추천 23개 · 혼잡 주의 21개 — ", {}),
         ("종로1.2.3.4가동이 적합도 98.5로 1위", {"bold": True, "color": INK})],
        [("신촌동은 음식점 1,132개에 공영주차 ", {}), ("46면", {"bold": True, "color": RED}),
         (", 논현2동은 927개에 ", {}), ("13면", {"bold": True, "color": RED})],
    ], y=5.75)


def s_falsepos(prs, n):
    s = _content(prs, n, "분석 결과", "가설 6 — 아파트가 만드는 판정 오탐",
                 sub="공영주차 실수요는 부설주차장이 없는 가구에서 발생한다")
    _picture(s, "06_false_positive", top=2.5, height=3.05)
    _bullets(s, [
        [("비아파트가 ", {}), ("2~35호뿐인 아파트 단지", {"bold": True, "color": RED}),
         ("가 '공급부족'으로 지목되고 있었다", {})],
        [("거처 종류를 함께 보면 ", {}), ("확충 불필요 13개 동", {"bold": True, "color": INK}),
         ("을 예산 배분에서 먼저 제외할 수 있다", {})],
    ], y=5.75)


def s_priority(prs, n):
    s = _content(prs, n, "분석 결과", "확충 우선순위",
                 sub="순위는 회귀 잔차가 아니라 직접 지표로 냈다")
    _picture(s, "07_priority", top=2.4, height=3.25)
    _bullets(s, [
        [("설명력이 낮으면 잔차는 '보정한 값'이 아니라 ", {}),
         ("주차면 수 그 자체", {"bold": True, "color": INK}), ("가 된다 (순위 상관 0.997)", {})],
        [("보정한 척하지 않기 위해 ", {}),
         ("'비아파트 100가구당 주차면'", {"bold": True, "color": INK}), ("으로 제시했다", {})],
    ], y=5.8)


# ── IV. 결론 ──────────────────────────────────────────────────
def s_conclusion_user(prs, n):
    s = _content(prs, n, "결론", "나들이 이용자에게")
    _table(s, [
        ["질문", "답 (토요일 오후 기준)"],
        ["어디로 갈까", "종로1.2.3.4가동 — 적합도 98.5 · 음식점 2,047개 · 공영주차 1,595면"],
        ["언제 갈까", "을지로동 — 화·점심보다 일·밤이 379% 여유"],
        ["어디를 피할까", "신촌동 — 음식점 1,132개에 공영주차 46면"],
    ], y=2.4, col_w=[2.6, 8.9], size=14, row_h=0.5)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M_LEFT), Inches(4.6),
                             Inches(11.5), Inches(1.5))
    box.adjustments[0] = 0.1
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write(tf, [[("목적지를 바꾸지 않고 시점만 옮겨도", {"size": 16, "color": WHITE})],
                [("85.2%의 동에서 20% 이상 주차가 여유해진다",
                  {"size": 26, "bold": True, "color": WHITE})]],
           align=PP_ALIGN.CENTER, space=4)


def s_conclusion_policy(prs, n):
    s = _content(prs, n, "결론", "주차 정책 담당자에게")
    items = [
        ("현재 배치는 수요 기반이 아니다", "유동수요로도 정주 실수요로도 설명되지 않는다 (R² < 0.07)"),
        ("생활인구만으로 대상지를 고르면 아파트 단지를 지목한다", "비아파트 2~35호뿐인 4개 동이 실제로 그렇게 잡혔다"),
        ("확충 불필요 13개 동을 예산 배분에서 먼저 제외하라", "송파 5 · 노원 4 · 서초 2 · 강동 1 · 용산 1"),
        ("방문객 수요형은 거주자우선이 아닌 시간제로 풀어야 한다", "신촌동·화양동은 아파트가 9~11%뿐인데 방문객 수요가 압도적이다"),
    ]
    y = 2.4
    for i, (title, desc) in enumerate(items, 1):
        _write(_tf(s, M_LEFT, y, 0.5, 0.5), [f"{i:02d}"], size=17, bold=True, color=BLUE, space=0)
        _write(_tf(s, M_LEFT + 0.6, y - 0.04, 11.0, 0.9),
               [[(title, {"size": 17, "bold": True, "color": INK})],
                [(desc, {"size": 13, "color": MUTED})]], space=2)
        y += 1.08


def s_demo(prs, n):
    s = _content(prs, n, "결론", "대시보드 시연",
                 sub="모든 메뉴가 아니라 나들이객 한 명이 실제로 쓰는 흐름을 보여드립니다")
    flow = [("오늘 어디 갈까", "탐색 — 요일·시간대를 고르면 서울 전체 여유도가 바뀐다"),
            ("이 동네 괜찮나", "지도 — 적합도 · 확충 후보 · 공영주차 0면 동"),
            ("언제 가지", "코스 추천 — 동네별 최적 시점"),
            ("정책 담당자 화면", "확충 우선순위 — 자치구별 후보와 불필요 지역")]
    y = 2.5
    for i, (title, desc) in enumerate(flow, 1):
        _pill(s, M_LEFT, y, f"{i}. {title}", fill=BLUE if i < 4 else DARK, size=12.5)
        _write(_tf(s, M_LEFT + 3.4, y - 0.02, 8.5, 0.5), [desc], size=14.5, color=BODY, space=0)
        y += 0.82
    _hairline(s, 6.05)
    _write(_tf(s, M_LEFT, 6.2, 11.5, 0.4),
           [("streamlit run dashboard/app.py  →  localhost:8501", {"size": 13, "color": MUTED})],
           space=0)


def s_limits(prs, n):
    s = _content(prs, n, "결론", "한계와 향후 과제")
    _table(s, [
        ["한계", "내용"],
        ["잔차 ≠ 수요 대비 부족", "설명력이 낮아 잔차 순위가 주차면수 순위와 0.997 일치. 확충 순위는 직접 지표로 제시"],
        ["민영주차장 미확보", "표준데이터 수집범위가 '지자체 관리 대상'이라 서울 25개 구 중 2개 구 45건만 공개"],
        ["생활인구는 대리변수", "이동통신 기반 추계라 차량 이용자를 구분할 수 없다"],
        ["아파트 = 주차 해결 아님", "1990년대 이전 아파트는 세대당 0.3~0.7대로 현 기준 미달"],
    ], y=2.35, col_w=[3.0, 8.5], size=13, row_h=0.46)
    # 표 아래 4.65"부터 시작 — 표 높이를 고정했으므로 이 좌표가 어긋나지 않는다
    _write(_tf(s, M_LEFT, 4.85, 11.5, 0.42),
           [("향후 과제", {"size": 17, "bold": True, "color": INK})], space=0)
    _bullets(s, [
        "자치구별 자동차 등록대수를 결합해 '차량 보유 가구' 기준으로 정밀화",
        "주차요금을 결합해 '공영 부족 → 민영 비용 부담' 경로를 정량화",
        "동일 분석 틀을 계절 단위·타 광역시로 확장",
    ], y=5.35, size=14)


def s_learned(prs, n):
    s = _content(prs, n, "결론", "가장 크게 배운 점")
    items = [
        ("데이터가 안 붙는 것이 분석의 절반이었다",
         "코드 체계 불일치로 조인율이 7.5%였다. 크로스워크를 만들어 100%로 올리는 데 가장 많은 시간이 들었고,\n이 단계가 뒤의 모든 결과를 좌우했다"),
        ("가설이 기각되는 것도 결과다",
         "R²=0.068은 처음엔 실패로 보였다. 그러나 '공급이 수요를 따라가지 않는다'는 사실 자체가\n정책적으로 가장 쓸모 있는 발견이었다"),
        ("보정한 척하지 않는 것이 중요했다",
         "설명력이 낮으면 잔차는 보정된 값이 아니다. 이를 확인하고 순위를 직접 지표로 바꾼 것이\n분석의 정직성을 지켰다"),
    ]
    y = 2.35
    for title, desc in items:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(M_LEFT), Inches(y),
                                 Inches(0.06), Inches(1.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        bar.shadow.inherit = False
        _write(_tf(s, M_LEFT + 0.22, y - 0.05, 11.3, 1.2),
               [[(title, {"size": 17, "bold": True, "color": INK})],
                [(desc, {"size": 13, "color": BODY})]], space=3)
        y += 1.42
    _note(s, "산출물 — 보고서 reports/report.md · 그림 19종 · 데이터 8종 · 대시보드 5페이지", y=6.75)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    plan = [
        (s_cover, None), (s_agenda, None),
        (s_divider, ("01", "분석 배경 및 목표")),
        (s_scope, None), (s_problem, None), (s_hypotheses, None), (s_pipeline, None),
        (s_divider, ("02", "활용 데이터")),
        (s_data, None), (s_join, None), (s_timeslot_design, None),
        (s_divider, ("03", "분석 결과")),
        (s_eda, None),
        (s_core, None), (s_hyp13, None), (s_cluster, None), (s_timeslot, None),
        (s_recommend, None), (s_falsepos, None), (s_priority, None),
        (s_divider, ("04", "결론")),
        (s_conclusion_user, None), (s_conclusion_policy, None), (s_demo, None),
        (s_limits, None), (s_learned, None),
    ]
    for i, (build, args) in enumerate(plan, 1):
        build(prs, i, *args) if args else build(prs, i)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    logger.info(f"저장 완료: {OUT_PATH} ({len(prs.slides)}장)")
    print(f"\n{len(prs.slides)}장 생성 -> {OUT_PATH}")


if __name__ == "__main__":
    main()
