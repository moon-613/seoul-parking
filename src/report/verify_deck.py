"""발표자료에 적힌 숫자가 데이터와 맞는지 검사 -> 어긋나면 exit 1

왜 필요한가
----------
`deck_figures.py`가 그리는 **그림**은 매번 데이터를 읽어 다시 계산하지만,
`build_deck.py`의 **표와 문장**은 손으로 적은 문자열이다.
분석을 다시 돌려 수치가 바뀌어도 슬라이드는 따라오지 않는다.

실제로 이 방식 때문에 사고가 있었다(PROGRESS ⑱).
핵심 결과 슬라이드에 정주수요 R²를 0.004로 적어뒀는데 그림은 0.018짜리 표본을
그리고 있었다. 사람이 대조해서 겨우 찾았지, 구조가 막아준 것이 아니었다.

표를 전부 자동 생성하는 방법도 있지만 서술형 표가 절반이라 이득이 크지 않다.
대신 **적힌 값과 계산한 값을 대조**해 어긋나면 알려주는 쪽을 택했다.
코드 복잡도는 그대로 두면서 사고만 막는다.

어떻게 검사하나
-------------
1. 생성된 .pptx에서 모든 텍스트(도형·표)를 뽑는다
2. 각 항목마다 **정규식으로 슬라이드에 적힌 값**을 찾고
3. 같은 값을 **data/processed/ 에서 직접 계산**해 비교한다

  발견 못 함  -> 경고 (문구를 바꿨거나 슬라이드에서 뺐다는 뜻)
  값이 다름   -> 실패 (데이터가 바뀌었는데 슬라이드가 안 따라왔다)

실행
----
  python -m src.report.verify_deck
  python -m src.report.verify_deck --strict   (경고도 실패로 처리)
"""
from __future__ import annotations

import argparse
import re
import sys

import numpy as np
import pandas as pd
from pptx import Presentation
from scipy import stats

from src.utils.logger import get_logger
from src.utils.settings import DATA_PROCESSED, ROOT_DIR

logger = get_logger(__name__)

DECK_PATH = ROOT_DIR / "submission" / "02_발표자료" / "분석보고서_문지영.pptx"
MIN_NONAPT = 100          # real_demand.py 와 같은 값
BASE = ("토", "오후")      # config 의 regression_baseline


# ── 데이터에서 값 계산 ──────────────────────────────────────────
def _load(name: str) -> pd.DataFrame:
    return pd.read_csv(DATA_PROCESSED / name, dtype={"adm_cd": str, "admi_cd": str})


def facts() -> dict[str, float]:
    """슬라이드가 주장하는 값을 데이터에서 다시 계산한다."""
    panel = _load("panel.csv")
    base = panel[(panel.weekday == BASE[0]) & (panel.timeslot == BASE[1])]
    parked = base[base.has_parking]
    rec = _load("dong_recommend.csv")
    rd = _load("dong_real_demand.csv")
    ts = _load("dong_timeslot.csv")
    cl = _load("dong_cluster.csv")
    suit = _load("dong_suitability.csv")
    sgg = pd.read_csv(DATA_PROCESSED / "sgg_summary.csv")

    def r2(x, y):
        return stats.linregress(np.log1p(x), np.log1p(y)).rvalue ** 2

    def corr(col):
        ok = parked.dropna(subset=[col, "slots_per_1k"])
        return stats.pearsonr(ok[col], ok.slots_per_1k)[0]

    flow = base.dropna(subset=["non_apt_households", "living_pop", "parking_slots"])
    flow = flow[flow.has_parking]
    res = flow[flow.non_apt_households >= MIN_NONAPT]
    top = suit.nlargest(1, "나들이적합도").iloc[0]

    def dong(df, name, col):
        return df.loc[df.admi_nm == name, col].iloc[0]

    return {
        # 패널
        "패널 행수": len(panel),
        "패널 열수": panel.shape[1],
        "행정동 수": panel.adm_cd.nunique(),
        "시간대 수": panel.timeslot.nunique(),
        # 공급
        "총 공영주차면": parked.parking_slots.sum(),
        "주차장 보유 동": parked.adm_cd.nunique(),
        "0면 동": (~base.has_parking).sum(),
        # 회귀
        "R2 유동수요": r2(flow.living_pop, flow.parking_slots),
        "R2 정주수요": r2(res.non_apt_households, res.parking_slots),
        "상관 생활인구": corr("living_pop"),
        "상관 음식점": corr("store_food"),
        # 등급
        "추천 동": (rec.등급 == "추천").sum(),
        "혼잡 주의 동": (rec.등급 == "혼잡 주의").sum(),
        "확충 후보": sgg.확충후보.sum(),
        "확충 불필요": (rd.확충필요성 == "불필요").sum(),
        # 군집
        "상권형 동": (cl.유형 == "상권형").sum(),
        "주거형 동": (cl.유형 == "주거형").sum(),
        # 시간대 효과
        "개선율 중앙값": ts.개선율.median() * 100,
        "개선율 20퍼센트 이상 비율": (ts.개선율 >= 0.20).mean() * 100,
        "소공동 개선율": dong(ts, "소공동", "개선율") * 100,
        "을지로동 개선율": dong(ts, "을지로동", "개선율") * 100,
        # 개별 동
        "적합도 1위 점수": top.나들이적합도,
        "신촌동 음식점": dong(rec, "신촌동", "store_food"),
        "신촌동 주차면": dong(rec, "신촌동", "parking_slots"),
        "논현2동 음식점": dong(rec, "논현2동", "store_food"),
        "논현2동 주차면": dong(rec, "논현2동", "parking_slots"),
        # 자치구
        "자치구 중앙값": sgg.실수요100가구당.median(),
    }


# ── 슬라이드에서 값 추출 ────────────────────────────────────────
def deck_text() -> str:
    if not DECK_PATH.exists():
        raise FileNotFoundError(
            f"발표자료가 없습니다: {DECK_PATH}\n"
            "먼저 python -m src.report.build_deck 를 실행하세요."
        )
    prs = Presentation(DECK_PATH)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(c.text for c in row.cells)
    # 줄바꿈·연속 공백을 하나로 눌러 정규식이 줄에 걸리지 않게 한다
    return re.sub(r"\s+", " ", " ".join(parts))


# 항목: (이름, 슬라이드에서 찾을 정규식, 허용 오차)
#   정규식의 첫 그룹이 비교할 숫자다. 쉼표는 제거하고 float으로 읽는다.
CHECKS: list[tuple[str, str, float]] = [
    ("패널 행수", r"([\d,]+)행 패널", 0),
    ("행정동 수", r"행정동 ([\d,]+) × 요일", 0),
    ("시간대 수", r"요일 7 × 시간대 (\d+)", 0),
    ("총 공영주차면", r"([\d,]+)면 / \d+개 동", 0),
    ("주차장 보유 동", r"[\d,]+면 / (\d+)개 동", 0),
    ("0면 동", r"공영주차 0면 (\d+)개 제외", 0),
    # 두 곳(4번 목표·18번 가설4)에 나오므로 둘 다 잡아 같은 값인지 본다
    ("추천 동", r"추천 동네 (\d+)개 · 확충 후보", 0),
    ("추천 동", r"• 추천 (\d+)개 · 혼잡 주의", 0),
    ("확충 후보", r"확충 후보 (\d+)개", 0),
    ("확충 불필요", r"확충 불필요 (\d+)개", 0),
    ("혼잡 주의 동", r"혼잡 주의 (\d+)개", 0),
    ("상권형 동", r"상권형 (\d+)개 동", 0),
    ("주거형 동", r"주거형 (\d+)개 동", 0),
    ("개선율 20퍼센트 이상 비율", r"([\d.]+)%의 동에서 20% 이상", 0.05),
    ("소공동 개선율", r"소공동은 수·점심 대비 일·밤이 ([\d,]+)%", 0.5),
    ("을지로동 개선율", r"을지로동 — 화·점심보다 일·밤이 ([\d,]+)%", 0.5),
    ("적합도 1위 점수", r"적합도 ([\d.]+) ·", 0.05),
    ("신촌동 음식점", r"신촌동 — 음식점 ([\d,]+)개", 0),
    ("신촌동 주차면", r"신촌동 — 음식점 [\d,]+개에 공영주차 (\d+)면", 0),
    ("논현2동 음식점", r"논현2동은 ([\d,]+)개에", 0),
    ("논현2동 주차면", r"논현2동은 [\d,]+개에 (\d+)면", 0),
    ("상관 생활인구", r"r = (-[\d.]+)\s*$|약한 지지\s+r = (-[\d.]+)", 0.005),
]


def run(strict: bool = False) -> int:
    text = deck_text()
    data = facts()

    ok = warn = fail = 0
    logger.info(f"검사 대상 {DECK_PATH.name} · 항목 {len(CHECKS)}개")
    for name, pattern, tol in CHECKS:
        m = re.search(pattern, text)
        if not m:
            logger.warning(f"  [찾음X] {name:24} 슬라이드에서 문구를 찾지 못했습니다")
            warn += 1
            continue
        raw = next(g for g in m.groups() if g)
        said = float(raw.replace(",", ""))
        real = float(data[name])
        if abs(said - real) <= tol:
            logger.info(f"  [일치 ] {name:24} {said:,.10g}")
            ok += 1
        else:
            logger.error(f"  [불일치] {name:24} 슬라이드 {said:,.10g}  vs  데이터 {real:,.10g}")
            fail += 1

    # 회귀 수치는 그림이 계산해 넣으므로 값만 기록해 둔다 (사람이 눈으로 확인)
    logger.info("참고 — 그림이 직접 계산하는 값")
    for k in ["R2 유동수요", "R2 정주수요", "상관 음식점", "개선율 중앙값", "자치구 중앙값"]:
        logger.info(f"  {k:24} {data[k]:.3f}")

    logger.info(f"결과 — 일치 {ok} · 경고 {warn} · 불일치 {fail}")
    if fail or (strict and warn):
        logger.error("발표자료의 숫자가 데이터와 어긋납니다. build_deck.py를 고치세요.")
        return 1
    logger.info("발표자료의 숫자가 데이터와 일치합니다.")
    return 0


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--strict", action="store_true", help="문구를 못 찾은 경우도 실패로 처리")
    sys.exit(run(ap.parse_args().strict))


if __name__ == "__main__":
    main()
